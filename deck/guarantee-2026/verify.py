#!/usr/bin/env python3
"""생성한 PPTX 를 구조·원문·차트 세 축으로 검증한다.

원본에 텍스트 레이어가 없어 전사가 유일한 원문 경로다. content.json 의 모든
문자열이 실제로 PPTX 안에 들어갔는지 기계로 대조한다.

이 덱은 그래프를 **네이티브 차트**로 넣는다. 차트 안 글자는 도형 텍스트로 잡히지
않으므로 차트의 항목·값까지 따로 훑어 대조 대상에 넣고, 편집용 워크북이 실제로
들어 있는지도 확인한다(없으면 파워포인트에서 '데이터 편집'이 열리지 않는다).

사용: python3 verify.py [pptx경로] [--ids g01,g02]
"""
import json
import re
import sys
import unicodedata
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

ROOT = Path(__file__).parent
EXPECT_W_IN = 10.667
EXPECT_H_IN = 16.0
MIN_SHAPES = 20        # kstat-ppt 원칙 17 — 도형 밀도
MIN_PT = 10            # kstat-ppt 원칙 15 — 본문 최소 폰트

FORBIDDEN_PRST = {'roundRect', 'round1Rect', 'round2SameRect', 'round2DiagRect'}

# 장별 기대 차트 수 (레이아웃을 고칠 때 함께 고친다)
EXPECT_CHARTS = {'g01': 1, 'g02': 2, 'g03': 8, 'g04': 7}
# 장별 기대 이미지 수
EXPECT_PICS = {'g01': 16, 'g02': 24, 'g03': 25, 'g04': 35}


def norm(s):
    s = unicodedata.normalize('NFC', str(s))
    s = s.replace('‘', "'").replace('’', "'").replace('“', '"').replace('”', '"')
    s = s.replace('‧', '·')
    return re.sub(r'\s+', '', s)


def walk_strings(node, out):
    """content.json 에서 화면에 찍히는 문자열만 모은다."""
    if isinstance(node, str):
        if node.strip() and not re.fullmatch(r'g\d\d-[a-z0-9-]+', node):
            out.append(node)
    elif isinstance(node, list):
        for v in node:
            walk_strings(v, out)
    elif isinstance(node, dict):
        for k, v in node.items():
            if k in ('id', 'img', 'art', 'icon', 'noteIcon', 'fixes', 'cats'):
                continue
            walk_strings(v, out)


def shape_text(shape, bag, charts):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sp in shape.shapes:
            shape_text(sp, bag, charts)
        return
    if getattr(shape, 'has_chart', False) and shape.has_chart:
        ch = shape.chart
        info = {'cats': [], 'vals': []}
        try:
            info['cats'] = [str(c) for c in ch.plots[0].categories]
        except Exception:
            pass
        for se in ch.series:
            info['vals'].extend([v for v in se.values])
        charts.append(info)
        bag.extend(info['cats'])
        return
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            bag.append(''.join(r.text for r in p.runs))


def fonts_and_sizes(shape, fonts, small):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sp in shape.shapes:
            fonts_and_sizes(sp, fonts, small)
        return
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if r.font.name:
                fonts.add(r.font.name)
            if r.font.size and r.font.size.pt < MIN_PT - 0.01 and r.text.strip():
                small.append((round(r.font.size.pt, 1), r.text[:20]))


def main():
    args = [a for a in sys.argv[1:] if not a.startswith('--')]
    ids = None
    if '--ids' in sys.argv:
        ids = sys.argv[sys.argv.index('--ids') + 1].split(',')

    if args:
        path = Path(args[0])
    else:
        cand = sorted((ROOT / 'out').glob('*.pptx'), key=lambda p: p.stat().st_mtime)
        if not cand:
            sys.exit('out/ 에 pptx 가 없다')
        path = cand[-1]

    content = json.loads((ROOT / 'content.json').read_text())
    if ids:
        content = [d for d in content if d['id'] in ids]

    prs = Presentation(str(path))
    fails = []
    print(f'파일: {path}')

    w_in = Emu(prs.slide_width).inches
    h_in = Emu(prs.slide_height).inches
    if abs(w_in - EXPECT_W_IN) > 0.02 or abs(h_in - EXPECT_H_IN) > 0.02:
        fails.append(f'슬라이드 크기 {w_in:.3f}×{h_in:.3f}in != {EXPECT_W_IN}×{EXPECT_H_IN}')

    if len(prs.slides) != len(content):
        fails.append(f'슬라이드 수 {len(prs.slides)} != 기대 {len(content)}')

    all_charts = []
    for i, (slide, d) in enumerate(zip(prs.slides, content), 1):
        bag, charts, fonts, small = [], [], set(), []
        pics = 0
        for sh in slide.shapes:
            shape_text(sh, bag, charts)
            fonts_and_sizes(sh, fonts, small)
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
            prst = getattr(getattr(sh, 'element', None), 'find', None)
            if sh.has_text_frame is not None:
                pass
        all_charts.append((d['id'], charts))

        n = len(slide.shapes)
        print(f'  {d["id"]}: 도형 {n}개 / 이미지 {pics}개 / 차트 {len(charts)}개')
        if n < MIN_SHAPES:
            fails.append(f'{d["id"]}: 도형 {n}개 < {MIN_SHAPES} (밀도 부족)')
        if d['id'] in EXPECT_CHARTS and len(charts) != EXPECT_CHARTS[d['id']]:
            fails.append(f'{d["id"]}: 차트 {len(charts)}개 != 기대 {EXPECT_CHARTS[d["id"]]}개')
        if d['id'] in EXPECT_PICS and pics != EXPECT_PICS[d['id']]:
            fails.append(f'{d["id"]}: 이미지 {pics}개 != 기대 {EXPECT_PICS[d["id"]]}개')
        for pt, t in small:
            fails.append(f'{d["id"]}: {pt}pt < {MIN_PT}pt — "{t}"')

        # 금지 도형
        xml = slide.shapes._spTree.xml
        for prst in FORBIDDEN_PRST:
            if f'prst="{prst}"' in xml:
                fails.append(f'{d["id"]}: 금지 도형 {prst} 사용')

        # 원문 대조
        blob = norm(''.join(bag))
        want = []
        walk_strings(d, want)
        miss = [t for t in want if norm(t) and norm(t) not in blob]
        if miss:
            for t in miss:
                fails.append(f'{d["id"]} 원문 누락: {t}')
        print(f'      원문 조각 {len(want)}개 중 누락 {len(miss)}건')

        # 차트 값 대조 — content.json 의 vals 와 자릿수까지 같아야 한다
        want_vals = []
        def collect(node):
            if isinstance(node, dict):
                if 'vals' in node and isinstance(node['vals'], list):
                    want_vals.append([float(v) for v in node['vals']])
                for v in node.values():
                    collect(v)
            elif isinstance(node, list):
                for v in node:
                    collect(v)
        collect(d)
        got_vals = [[round(float(v), 6) for v in c['vals']] for c in charts]
        # 가로 막대는 첫 항목을 위에 두려고 배열을 뒤집어 넣는다(charts.js).
        # 그래서 뒤집힌 순서도 일치로 본다.
        for wv in want_vals:
            def same(gv, wv=wv):
                if len(gv) != len(wv):
                    return False
                return (all(abs(a - b) < 1e-6 for a, b in zip(gv, wv))
                        or all(abs(a - b) < 1e-6 for a, b in zip(gv, wv[::-1])))
            if not any(same(gv) for gv in got_vals):
                fails.append(f'{d["id"]}: 차트 값 {wv} 이(가) 산출물에 없다')

    # 편집용 워크북 — 없으면 '데이터 편집'이 열리지 않는다
    with zipfile.ZipFile(path) as z:
        books = [n for n in z.namelist() if n.startswith('ppt/embeddings/')]
        chart_xml = [n for n in z.namelist() if re.match(r'ppt/charts/chart\d+\.xml$', n)]
    print(f'  차트 XML {len(chart_xml)}개 / 편집용 워크북 {len(books)}개')
    if len(books) < len(chart_xml):
        fails.append(f'편집용 워크북 {len(books)}개 < 차트 {len(chart_xml)}개 '
                     '— 파워포인트에서 데이터 편집이 열리지 않는다')

    # 절대 경로 유출 (pptxgenjs 가 altText 를 안 주면 소스 경로를 넣는다)
    with zipfile.ZipFile(path) as z:
        for n in z.namelist():
            if n.startswith('ppt/slides/slide') and n.endswith('.xml'):
                if '/home/' in z.read(n).decode('utf8'):
                    fails.append(f'{n}: 빌드한 사람의 절대 경로가 남아 있다')

    if fails:
        print(f'\n실패 {len(fails)}건:')
        for f in fails[:60]:
            print(f'  - {f}')
        sys.exit(1)
    print('\n전 항목 통과')


main()
