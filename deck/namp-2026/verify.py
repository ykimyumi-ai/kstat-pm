#!/usr/bin/env python3
"""생성한 PPTX를 구조·원문 두 축으로 검증한다.

원본 PDF에 텍스트 레이어가 없어 전사가 유일한 원문 경로이므로,
content.js 의 모든 문자열이 PPTX 안에 실제로 들어갔는지 기계로 대조한다.

사용: python3 verify.py [pptx경로] [--ids s12,s13]
"""
import json
import re
import subprocess
import sys
import unicodedata
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

ROOT = Path(__file__).parent
EXPECT_W_IN = 11.0
EXPECT_H_IN = 7.3333
MIN_SHAPES = 20        # kstat-ppt 원칙 17 — 도형 밀도
MIN_PT = 10            # kstat-ppt 원칙 15 — 본문 최소 폰트

FORBIDDEN_PRST = {'roundRect', 'round1Rect', 'round2SameRect', 'round2DiagRect'}


def norm(s):
    """대조용 정규화 — 공백·유사문자 차이를 흡수한다."""
    s = unicodedata.normalize('NFC', s)
    s = s.replace('‘', "'").replace('’', "'")
    s = s.replace('“', '"').replace('”', '"')
    s = s.replace('·', '·').replace('‧', '·')
    s = re.sub(r'\s+', '', s)
    return s


def slide_text(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            out.append(sh.text_frame.text)
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    out.append(cell.text)
    return '\n'.join(out)


def collect_expected():
    """content.js 를 node로 읽어 슬라이드별 기대 문자열을 뽑는다."""
    js = r'''
      const c = require('./content');
      const out = c.map(d => {
        const acc = [];
        const walk = (v) => {
          if (v == null) return;
          if (typeof v === 'string') { acc.push(v); return; }
          if (Array.isArray(v)) { v.forEach(walk); return; }
          if (typeof v === 'object') {
            // 화면에 찍히는 글자가 아니라 서식·메타 표시자
            const META = new Set(['img', 'img2', 'icon', 'id', 'options', 'c',
                                  'accent', 'tone', 'tint', 'photo', 'quoteStyle', 'fixes']);
            for (const k of Object.keys(v)) {
              if (META.has(k)) continue;
              walk(v[k]);
            }
          }
        };
        walk(d);
        return { id: d.id, strings: acc };
      });
      process.stdout.write(JSON.stringify(out));
    '''
    r = subprocess.run(['node', '-e', js], cwd=ROOT, capture_output=True, text=True)
    r.check_returncode()
    return json.loads(r.stdout)


def main():
    args = [a for a in sys.argv[1:] if not a.startswith('--')]
    ids = None
    if '--ids' in sys.argv:
        ids = sys.argv[sys.argv.index('--ids') + 1].split(',')
    path = args[0] if args else str(sorted((ROOT / 'out').glob('*.pptx'))[0])
    prs = Presentation(path)
    fails, warns = [], []

    expected = collect_expected()
    if ids:
        expected = [e for e in expected if e['id'] in ids]

    # ── ① 구조 ────────────────────────────────────────────
    n = len(prs.slides)
    if n != len(expected):
        fails.append(f'슬라이드 수 {n} != 기대 {len(expected)}')
    w_in = Emu(prs.slide_width).inches
    h_in = Emu(prs.slide_height).inches
    if abs(w_in - EXPECT_W_IN) > 0.01 or abs(h_in - EXPECT_H_IN) > 0.01:
        fails.append(f'슬라이드 크기 {w_in:.3f}×{h_in:.3f} != {EXPECT_W_IN}×{EXPECT_H_IN}')

    xml_all = ''
    tables = 0
    pics = 0
    for i, sl in enumerate(prs.slides, 1):
        shapes = list(sl.shapes)
        if len(shapes) < MIN_SHAPES:
            fails.append(f'슬라이드 {i}: 도형 {len(shapes)}개 < {MIN_SHAPES}')
        for sh in shapes:
            if sh.has_table:
                tables += 1
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
        x = sl._element.xml
        xml_all += x
        for prst in FORBIDDEN_PRST:
            if f'prst="{prst}"' in x:
                fails.append(f'슬라이드 {i}: 금지 도형 {prst}')
        if '<a:effectLst><a:outerShdw' in x:
            fails.append(f'슬라이드 {i}: shadow 사용')
        for m in re.finditer(r'srgbClr val="([0-9A-Fa-f]+)"', x):
            if len(m.group(1)) != 6:
                fails.append(f'슬라이드 {i}: 잘못된 hex {m.group(1)}')
        for m in re.finditer(r'\ssz="(\d+)"', x):
            pt = int(m.group(1)) / 100
            if pt < MIN_PT:
                fails.append(f'슬라이드 {i}: 폰트 {pt}pt < {MIN_PT}pt')

    # 네이티브 표는 10장 비목별 산출표에만 있다. 10장이 대상일 때만 요구한다.
    if any(e['id'] == 's10' for e in expected) and tables < 1:
        fails.append('네이티브 표가 하나도 없음 (10장 비목별 산출표 필요)')
    # 14~16장은 원본에 사진·아이콘이 있어 삽입 이미지가 반드시 있어야 한다.
    need_pic = {'s14': 9, 's15': 2, 's16': 3, 's17': 1}
    want = sum(v for k, v in need_pic.items() if any(e['id'] == k for e in expected))
    if want and pics < want:
        fails.append(f'삽입 이미지 {pics}개 < 기대 {want}개 (사진·아이콘 누락)')

    fonts = set(re.findall(r'typeface="([^"]+)"', xml_all))
    stray = {f for f in fonts if 'KoPub' not in f and f not in ('+mn-lt', '+mj-lt', 'Arial')}
    if stray:
        warns.append(f'KoPub 외 폰트 사용: {sorted(stray)}')

    # ── ② 원문 대조 ───────────────────────────────────────
    missing_total = 0
    for i, exp in enumerate(expected):
        if i >= n:
            break
        have = norm(slide_text(prs.slides[i]))
        for sstr in exp['strings']:
            for part in str(sstr).split('\n'):
                p = norm(part)
                if len(p) < 2:
                    continue
                if p not in have:
                    missing_total += 1
                    fails.append(f"슬라이드 {i+1}({exp['id']}) 원문 누락: {part[:48]}")

    # ── 결과 ─────────────────────────────────────────────
    print(f'파일: {path}')
    print(f'슬라이드 {n}장 / {w_in:.3f}×{h_in:.3f}in / 네이티브 표 {tables}개 / 삽입 이미지 {pics}개')
    print(f'슬라이드별 도형 수: '
          + ', '.join(str(len(list(s.shapes))) for s in prs.slides))
    print(f'사용 폰트: {sorted(fonts)}')
    print(f'원문 조각 대조: 누락 {missing_total}건')
    for w in warns:
        print(f'  [경고] {w}')
    if fails:
        print(f'\n실패 {len(fails)}건:')
        for f in fails[:60]:
            print(f'  - {f}')
        sys.exit(1)
    print('\n전 항목 통과')


if __name__ == '__main__':
    main()
